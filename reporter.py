"""
reporter.py — генерация DOCX и PPTX отчётов из брифинга.
Использует python-docx и python-pptx (бесплатно, без платных сервисов).
"""

import logging
from datetime import date
from pathlib import Path
from analyzer import run_all_scenarios, check_triggers, ScenarioResult
from collector import get_latest

log = logging.getLogger(__name__)

REPORTS_DIR = Path(__file__).parent / "reports"
REPORTS_DIR.mkdir(exist_ok=True)


# ---------------------------------------------------------------------------
# DOCX — ежедневный брифинг
# ---------------------------------------------------------------------------

def save_docx(briefing_text: str, filename: str | None = None) -> Path:
    """Сохраняет брифинг в .docx."""
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    today = date.today().isoformat()
    out_path = REPORTS_DIR / (filename or f"briefing_{today}.docx")

    doc = Document()

    # Заголовок
    title = doc.add_heading(f"КМГ Стратегический брифинг — {today}", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Текст брифинга (разбиваем по строкам)
    for line in briefing_text.split("\n"):
        line = line.strip()
        if not line:
            doc.add_paragraph()
            continue
        if line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.startswith("**") and line.endswith("**"):
            p = doc.add_paragraph()
            run = p.add_run(line.strip("**"))
            run.bold = True
        else:
            doc.add_paragraph(line)

    doc.save(out_path)
    log.info("DOCX сохранён: %s", out_path)
    return out_path


# ---------------------------------------------------------------------------
# PPTX — презентация сценариев для Совета директоров
# ---------------------------------------------------------------------------

def save_pptx(scenarios: dict[str, ScenarioResult],
              alerts: list,
              brent: float,
              kzt_usd: float,
              filename: str | None = None) -> Path:
    """Создаёт PPTX-презентацию с ключевыми показателями."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    today = date.today().isoformat()
    out_path = REPORTS_DIR / (filename or f"board_deck_{today}.pptx")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # пустой слайд

    # ------------------------------------------------------------------
    # Слайд 1 — Титульный
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _add_text(slide, f"КМГ — Стратегический брифинг", Inches(0.5), Inches(2.5),
              Inches(12), Inches(1.2), font_size=36, bold=True, color=(0x1a, 0x3a, 0x5c))
    _add_text(slide, f"Дата: {today} | Brent: ${brent:.2f} | KZT/USD: {kzt_usd:.0f}",
              Inches(0.5), Inches(3.8), Inches(12), Inches(0.6), font_size=18,
              color=(0x44, 0x44, 0x44))

    # ------------------------------------------------------------------
    # Слайд 2 — Сценарии 2026
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _add_text(slide, "Сценарии 2026 — финансовые прогнозы", Inches(0.5), Inches(0.3),
              Inches(12), Inches(0.7), font_size=24, bold=True, color=(0x1a, 0x3a, 0x5c))

    s_data = [
        ("Оптимистичный (20%)", scenarios["optimistic"], (0x1a, 0x7a, 0x3c)),
        ("Базовый (55%)",       scenarios["base"],       (0x1a, 0x3a, 0x5c)),
        ("Стрессовый (25%)",    scenarios["stress"],     (0xc0, 0x39, 0x2b)),
    ]
    y_start = Inches(1.2)
    for i, (label, s, color) in enumerate(s_data):
        y = y_start + Inches(i * 1.8)
        _add_text(slide, label, Inches(0.5), y, Inches(3), Inches(0.4),
                  font_size=14, bold=True, color=color)
        metrics = (
            f"Brent ${s.brent:.0f}  |  "
            f"Выручка ${s.revenue:,} млн  |  "
            f"EBITDA ${s.ebitda:,} млн  |  "
            f"FCF ${s.fcf:,} млн  |  "
            f"Δ выручка {s.revenue_delta_pct:+.1f}%"
        )
        _add_text(slide, metrics, Inches(0.5), y + Inches(0.45), Inches(12.3), Inches(0.5),
                  font_size=13, color=(0x33, 0x33, 0x33))
        flags = []
        if s.budget_risk: flags.append("⚠ Дефицит бюджета РК")
        if s.rating_risk: flags.append("⚠ Риск рейтинга BBB−")
        if flags:
            _add_text(slide, "  ".join(flags), Inches(0.5), y + Inches(0.95),
                      Inches(12), Inches(0.4), font_size=11, color=(0xc0, 0x39, 0x2b))

    # ------------------------------------------------------------------
    # Слайд 3 — Активные риски и алерты
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _add_text(slide, "Активные риски", Inches(0.5), Inches(0.3),
              Inches(12), Inches(0.7), font_size=24, bold=True, color=(0x1a, 0x3a, 0x5c))

    level_colors = {
        "КРИТИЧЕСКИЙ": (0xc0, 0x39, 0x2b),
        "ВЫСОКИЙ":     (0xe6, 0x7e, 0x22),
        "ПЛАНОВЫЙ":    (0x27, 0x6e, 0x3a),
    }
    for i, a in enumerate(alerts[:6]):
        y = Inches(1.1) + Inches(i * 0.9)
        color = level_colors.get(a.level, (0x33, 0x33, 0x33))
        _add_text(slide, f"[{a.level}]", Inches(0.5), y, Inches(1.8), Inches(0.5),
                  font_size=12, bold=True, color=color)
        _add_text(slide, a.message, Inches(2.4), y, Inches(10.4), Inches(0.5),
                  font_size=11, color=(0x22, 0x22, 0x22))

    # ------------------------------------------------------------------
    # Слайд 4 — Рекомендации (7 решений из PRD §4)
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _add_text(slide, "Стратегические рекомендации — статус", Inches(0.5), Inches(0.3),
              Inches(12), Inches(0.7), font_size=24, bold=True, color=(0x1a, 0x3a, 0x5c))

    recommendations = [
        ("1. Пересмотр бюджета 2026–2027",              "До 20 июня 2026"),
        ("2. Диверсификация маршрутов (КТК → 50–55%)",  "До 30 июня 2026"),
        ("3. Хеджирование 20–25% объёма Brent",         "Немедленно"),
        ("4. Переговоры ОПЕК+ — статус Тенгиза",        "Июль 2026"),
        ("5. Арбитраж Кашаган — мировое соглашение",    "В процессе"),
        ("6. ESG-стратегия + зелёные облигации",         "Июнь 2026"),
        ("7. Петрохимия + СПГ к 2030",                  "Горизонт 2028–2030"),
    ]
    for i, (rec, deadline) in enumerate(recommendations):
        y = Inches(1.1) + Inches(i * 0.83)
        _add_text(slide, rec,      Inches(0.5),  y, Inches(9.5), Inches(0.55), font_size=12)
        _add_text(slide, deadline, Inches(10.1), y, Inches(3.0), Inches(0.55), font_size=11,
                  color=(0x44, 0x44, 0x88))

    prs.save(out_path)
    log.info("PPTX сохранён: %s", out_path)
    return out_path


# ---------------------------------------------------------------------------
# Вспомогательная функция добавления текста
# ---------------------------------------------------------------------------

def _add_text(slide, text: str, left, top, width, height,
              font_size: int = 14, bold: bool = False,
              color: tuple = (0x00, 0x00, 0x00)):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


# ---------------------------------------------------------------------------
# Точка входа
# ---------------------------------------------------------------------------

def generate_reports(briefing_text: str | None = None) -> dict[str, Path]:
    brent   = get_latest("Brent_USD") or 85.0
    kzt_usd = get_latest("KZT_USD")   or 485.0
    scenarios = run_all_scenarios(brent)
    alerts    = check_triggers(brent, kzt_usd)

    paths = {}

    if briefing_text:
        paths["docx"] = save_docx(briefing_text)

    paths["pptx"] = save_pptx(scenarios, alerts, brent, kzt_usd)
    return paths


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
    result = generate_reports("Тестовый брифинг для проверки генерации документов.")
    for fmt, path in result.items():
        print(f"  {fmt.upper()}: {path}")
